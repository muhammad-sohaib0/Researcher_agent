# Tools Wrapper for API
# This module provides standalone functions for file reading without agent SDK dependencies

import os
from pathlib import Path
from dotenv import load_dotenv

# Load environment
load_dotenv(Path(__file__).parent.parent / ".env")


def read_pdf_tool(file_path: str) -> str:
    """
    Read and extract text from a PDF file using Groq Llama 4 Scout Vision OCR.
    Uses Llama 4 multimodal model to perform OCR on PDF pages.
    """
    try:
        from groq import Groq
        import fitz  # PyMuPDF for PDF to image conversion
        import base64
        
        # Check if file exists
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        # Get Groq API key
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "❌ Error: groq_api_key not found in environment variables"
        
        # Initialize Groq client
        client = Groq(api_key=groq_key)
        
        print(f"📄 Processing PDF with Groq Llama 4 Scout OCR: {Path(file_path).name}")
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(file_path)
        num_pages = len(pdf_document)
        
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        print(f"📦 PDF Size: {file_size_mb:.2f} MB, Pages: {num_pages}")
        
        # Extract text from all pages
        text_content = []
        text_content.append(f"📄 PDF Document: {Path(file_path).name}")
        text_content.append(f"📦 Size: {file_size_mb:.2f} MB")
        text_content.append(f"📊 Total Pages: {num_pages}")
        text_content.append(f"🤖 OCR Method: Groq Llama 4 Scout Vision")
        text_content.append("=" * 80)
        
        for page_num in range(num_pages):
            print(f"📖 Processing page {page_num + 1} of {num_pages}...")
            
            # Convert page to image
            page = pdf_document[page_num]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for quality
            
            # Convert to base64
            img_data = pix.tobytes("png")
            img_base64 = base64.b64encode(img_data).decode()
            
            # Send to Groq Llama 4 Scout (vision model)
            try:
                response = client.chat.completions.create(
                    model="meta-llama/llama-4-scout-17b-16e-instruct",
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Extract ALL text from this image. Perform accurate OCR. Preserve formatting and structure. Return ONLY the extracted text, nothing else."
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{img_base64}"
                                    }
                                }
                            ]
                        }
                    ],
                    max_tokens=4096
                )
                
                page_text = response.choices[0].message.content if response.choices else "[No text extracted]"
                
            except Exception as ocr_error:
                page_text = f"[OCR Error: {str(ocr_error)}]"
            
            # Add prominent page header
            text_content.append(f"\n\n{'#'*80}")
            text_content.append(f"#{'':^78}#")
            text_content.append(f"#{'📖 PAGE ' + str(page_num + 1) + ' OF ' + str(num_pages):^78}#")
            text_content.append(f"#{'':^78}#")
            text_content.append(f"{'#'*80}\n")
            text_content.append(page_text)
            text_content.append(f"\n{'─'*80}")
            text_content.append(f"[End of Page {page_num + 1}]")
            text_content.append(f"{'─'*80}")
        
        pdf_document.close()
        
        full_text = "\n".join(text_content)
        print(f"✅ Successfully processed all {num_pages} pages")
        return full_text
    
    except ImportError as ie:
        missing_lib = str(ie)
        if "groq" in missing_lib.lower():
            return "❌ Error: groq is required. Install it with: pip install groq"
        elif "fitz" in missing_lib or "PyMuPDF" in missing_lib:
            return "❌ Error: PyMuPDF is required. Install it with: pip install pymupdf"
        else:
            return f"❌ Error: Missing library - {missing_lib}"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"❌ Error reading PDF with OCR: {str(e)}"


def read_word_tool(file_path: str) -> str:
    """Read Word document and extract text."""
    try:
        from docx import Document
        
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        doc = Document(file_path)
        
        results = []
        results.append(f"📝 Word Document: {Path(file_path).name}")
        results.append("=" * 80)
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                results.append(paragraph.text)
        
        # Extract tables
        for table_num, table in enumerate(doc.tables):
            results.append(f"\nTable {table_num + 1}:")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                results.append(row_text)
        
        return "\n".join(results)
    except Exception as e:
        return f"Error reading Word document: {str(e)}"


def read_pptx_tool(file_path: str) -> str:
    """Read PowerPoint and extract text."""
    try:
        from pptx import Presentation
        
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        
        text_content = []
        text_content.append(f"📊 PowerPoint: {Path(file_path).name}")
        text_content.append(f"📑 Total Slides: {num_slides}")
        text_content.append("=" * 80)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(f"[Table Row] {row_text}")
            
            text_content.append(f"\n--- SLIDE {slide_num} ---")
            text_content.append("\n".join(slide_text))
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PowerPoint: {str(e)}"


def read_image_tool(file_path: str) -> str:
    """Read image and extract text via OCR using Groq Llama 4 Scout."""
    try:
        from groq import Groq
        import base64
        
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "❌ Error: groq_api_key not found"
        
        client = Groq(api_key=groq_key)
        
        # Read image
        with open(file_path, 'rb') as f:
            img_data = f.read()
        
        img_base64 = base64.b64encode(img_data).decode()
        
        # Determine mime type
        ext = Path(file_path).suffix.lower()
        mime_types = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.gif': 'image/gif', '.webp': 'image/webp'}
        mime_type = mime_types.get(ext, 'image/png')
        
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Extract ALL text from this image. Perform accurate OCR. Return ONLY the extracted text."},
                        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{img_base64}"}}
                    ]
                }
            ],
            max_tokens=4096
        )
        
        return response.choices[0].message.content if response.choices else "[No text extracted]"
    except Exception as e:
        return f"Error reading image: {str(e)}"


def extract_audio_tool(file_path: str, language: str = "en") -> str:
    """Extract text from audio file using Groq Whisper."""
    try:
        from groq import Groq
        
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "❌ Error: groq_api_key not found"
        
        client = Groq(api_key=groq_key)
        
        with open(file_path, "rb") as audio_file:
            transcription = client.audio.transcriptions.create(
                model="whisper-large-v3-turbo",
                file=audio_file,
                language=language,
                response_format="text"
            )
        
        return f"🎤 Transcription:\n{transcription}"
    except Exception as e:
        return f"Error transcribing audio: {str(e)}"
