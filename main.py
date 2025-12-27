import asyncio
from openai.types.responses import ResponseTextDeltaEvent
from openai import AsyncOpenAI
from agents import (Agent, Runner, GuardrailFunctionOutput, 
InputGuardrailTripwireTriggered, input_guardrail, AsyncOpenAI,
set_default_openai_client, set_tracing_disabled, set_default_openai_api, function_tool,
TResponseInputItem, ModelSettings, RunContextWrapper, ItemHelpers, trace, 
set_tracing_export_api_key, OpenAIChatCompletionsModel, handoff)
import os 
from dotenv import load_dotenv, find_dotenv
from pydantic import BaseModel
import PyPDF2
from pathlib import Path
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import pyttsx3
from serpapi import GoogleSearch
import aiohttp
import re
from urllib.parse import urlparse, unquote
import requests

# ONLY FOR TRACING
_: bool = load_dotenv(find_dotenv())
set_tracing_export_api_key(os.getenv("OPENAI_API_KEY", ""))


groq_key: str = os.getenv("groq_api_key", "")


# 1. Which LLM Service? - Groq
external_client: AsyncOpenAI = AsyncOpenAI(
    api_key=groq_key,
    base_url="https://api.groq.com/openai/v1",
)


# 2. Which LLM Model? - Kimi K2 Instruct
llm_model: OpenAIChatCompletionsModel = OpenAIChatCompletionsModel(
    model="moonshotai/kimi-k2-instruct",
    openai_client=external_client
)


# llmm_model: OpenAIChatCompletionsModel = OpenAIChatCompletionsModel(
#     model=="gemini-1.5-pro",
#     openai_client=external_client
# )



@function_tool
async def semantic_scholar_search(
    wrapper: RunContextWrapper, 
    query: str,
) -> str:
    """
    Search Semantic Scholar for research papers (12 results with open access PDFs).
    
    Args:
        query: The search query to use
        
    Returns:
        Structured list of PDF results from Semantic Scholar
    """
    try:
        print(f"[SEARCH] Searching Semantic Scholar: {query}")
        
        # Semantic Scholar API endpoint
        semantic_url = "https://api.semanticscholar.org/graph/v1/paper/search"
        semantic_params = {
            "query": query,
            "limit": 12,
            "fields": "title,abstract,url,openAccessPdf,year,authors,citationCount"
        }
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        
        # Optional: Add API key if you have one
        api_key = os.getenv("SEMANTIC_SCHOLAR_API_KEY")
        if api_key:
            headers["x-api-key"] = api_key
        
        semantic_response = requests.get(
            semantic_url, 
            params=semantic_params,
            headers=headers,
            timeout=15
        )
        
        if semantic_response.status_code != 200:
            return f"[ERROR] Semantic Scholar API error: {semantic_response.status_code}"
        
        semantic_data = semantic_response.json()
        pdf_results = []
        html_results = []
        
        for paper in semantic_data.get('data', []):
            # Get author names
            authors = ', '.join([
                author.get('name', '') 
                for author in paper.get('authors', [])[:3]
            ])
            if len(paper.get('authors', [])) > 3:
                authors += ' et al.'
            
            result_data = {
                'title': paper.get('title'),
                'snippet': paper.get('abstract', '')[:300] + '...' if paper.get('abstract') else '',
                'link': paper.get('url'),
                'source': 'Semantic Scholar',
                'year': paper.get('year'),
                'authors': authors,
                'citations': paper.get('citationCount', 0)
            }
            
            # Check for open access PDF
            if paper.get('openAccessPdf'):
                result_data['pdf_url'] = paper['openAccessPdf'].get('url')
                result_data['type'] = 'direct_pdf'
                pdf_results.append(result_data)
            else:
                result_data['type'] = 'html_page'
                html_results.append(result_data)
        
        all_results = pdf_results + html_results
        
        if not all_results:
            return "[ERROR] No results found from Semantic Scholar"

        print(f"[OK] Found {len(pdf_results)} open access PDFs and {len(html_results)} HTML pages")
        
        return {
            'total_results': len(all_results),
            'direct_pdfs': len(pdf_results),
            'html_pages': len(html_results),
            'source': 'Semantic Scholar',
            'results': all_results
        }
    
    except requests.exceptions.Timeout:
        return "[ERROR] Semantic Scholar search timeout"
    except Exception as e:
        return f"[ERROR] Semantic Scholar search error: {str(e)}"

@function_tool
async def google_scholar_search(
    wrapper: RunContextWrapper, 
    query: str,
) -> str:
    """
    Search Google Scholar for research papers (12 results with PDF prioritization).
    
    Args:
        query: The search query to use
        
    Returns:
        Structured list of PDF results from Google Scholar
    """
    try:
        # Add filetype:pdf to prioritize direct PDFs
        search_query = query + " filetype:pdf"
        
        params = {
            "engine": "google_scholar",
            "q": search_query,
            "num": "12",
            "api_key": os.getenv("SERPAPI_KEY")
        }
        
        print(f"🔍 Searching Google Scholar: {search_query}")
        search = GoogleSearch(params)
        results = search.get_dict()
        
        # Extract and prioritize direct PDF results
        pdf_results = []
        html_results = []
        
        for result in results.get('organic_results', []):
            result_data = {
                'title': result.get('title'),
                'snippet': result.get('snippet'),
                'link': result.get('link'),
                'source': 'Google Scholar'
            }
            
            # Check for direct PDF in resources
            if 'resources' in result:
                for resource in result.get('resources', []):
                    if resource.get('file_format') == 'PDF':
                        result_data['pdf_url'] = resource.get('link')
                        result_data['type'] = 'direct_pdf'
                        pdf_results.append(result_data)
                        break
            
            # If no direct PDF, it's an HTML result
            if 'pdf_url' not in result_data and result.get('link'):
                result_data['type'] = 'html_page'
                html_results.append(result_data)
        
        # Combine: direct PDFs first, then HTML pages
        all_results = pdf_results + html_results
        
        if not all_results:
            return "[ERROR] No results found from Google Scholar"
        
        print(f"[OK] Found {len(pdf_results)} PDFs and {len(html_results)} HTML pages")
        
        return {
            'total_results': len(all_results),
            'direct_pdfs': len(pdf_results),
            'html_pages': len(html_results),
            'source': 'Google Scholar',
            'results': all_results
        }
    
    except Exception as e:
        return f"[ERROR] Google Scholar search error: {str(e)}"

@function_tool
async def download_pdf(url: str, custom_filename: str = "") -> str:
    """
    Download a PDF file from a URL with improved reliability.
    - Extracts PDF from HTML pages
    - Adds security headers to bypass bot detection
    - Saves to project downloads folder for web access

    Args:
        url: The URL of the PDF file to download
        custom_filename: Optional custom filename (will auto-add .pdf extension)

    Returns:
        Success message with file path and download link
    """
    try:
        # Security headers to bypass bot detection
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Referer': 'https://scholar.google.com/',
            'Accept': 'application/pdf,*/*',
            'Accept-Language': 'en-US,en;q=0.9'
        }

        # Validate URL
        if not url.startswith(('http://', 'https://')):
            return "[ERROR] Error: URL must start with http:// or https://"

        # Use project downloads folder (for web serving)
        downloads_folder = Path(__file__).parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)

        print(f"[DOWNLOAD] Processing URL: {url}")
        
        # STEP 1: Check if this is an HTML page and extract PDF link
        if url.endswith('.html') or '/view' in url or '/article' in url:
            print(f"🔍 HTML page detected - extracting PDF link...")
            try:
                async with aiohttp.ClientSession(headers=headers) as session:
                    async with session.get(url, timeout=aiohttp.ClientTimeout(total=30), ssl=False) as response:
                        html_content = await response.text()
                        
                        # Search for PDF link patterns
                        pdf_patterns = [
                            r'href=["\']([^"\']*\.pdf[^"\']*)["\']',
                            r'src=["\']([^"\']*\.pdf[^"\']*)["\']',
                            r'data-pdf=["\']([^"\']*)["\']',
                            r'<a[^>]*href=["\']([^"\']*pdf[^"\']*)["\']',
                        ]
                        
                        pdf_url = None
                        for pattern in pdf_patterns:
                            matches = re.findall(pattern, html_content, re.IGNORECASE)
                            if matches:
                                pdf_candidate = matches[0]
                                if 'pdf' in pdf_candidate.lower():
                                    pdf_url = pdf_candidate
                                    break
                        
                        # If found, construct full URL
                        if pdf_url:
                            if not pdf_url.startswith('http'):
                                base_url = f"{urlparse(url).scheme}://{urlparse(url).netloc}"
                                pdf_url = base_url + pdf_url
                            
                            print(f"[OK] PDF link extracted: {pdf_url}")
                            url = pdf_url
                        else:
                            # Try alternate approach: replace /view with /pdf or /download
                            if '/view' in url:
                                alternate_url = url.replace('/view', '/pdf')
                                print(f"💡 Trying alternate URL: {alternate_url}")
                                url = alternate_url
            
            except Exception as e:
                print(f"[WARN] Warning extracting PDF from HTML: {e}")
                # Continue with original URL
        
        # STEP 2: Determine filename
        if custom_filename:
            filename = custom_filename if custom_filename.endswith('.pdf') else f"{custom_filename}.pdf"
        else:
            parsed_url = urlparse(url)
            filename = unquote(parsed_url.path.split('/')[-1])
            if not filename or not filename.endswith('.pdf'):
                filename = "downloaded_document.pdf"
        
        # Clean filename
        filename = re.sub(r'[<>:"|?*]', '_', filename)
        output_path = downloads_folder / filename
        
        print(f"💾 Saving to: {output_path}")
        
        # STEP 3: Download with improved headers and SSL bypass
        async with aiohttp.ClientSession(headers=headers) as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=300), ssl=False, allow_redirects=True) as response:
                
                # Check response status
                if response.status == 403:
                    return f"[ERROR] Access Denied (403): This site blocks automated downloads. Try manually downloading or find an alternative source."
                
                if response.status == 404:
                    return f"[ERROR] Not Found (404): The PDF URL is no longer available."
                
                if response.status != 200:
                    return f"[ERROR] Error: Failed to download. HTTP Status: {response.status}"
                
                # Check content type
                content_type = response.headers.get('Content-Type', '').lower()
                if 'pdf' not in content_type and not url.lower().endswith('.pdf'):
                    print(f"[WARN] Warning: Content type is '{content_type}', might not be PDF")
                
                # Get file size
                file_size = response.headers.get('Content-Length')
                if file_size:
                    size_mb = int(file_size) / (1024 * 1024)
                    print(f"📦 File size: {size_mb:.2f} MB")
                
                # STEP 4: Download with progress
                with open(output_path, 'wb') as f:
                    downloaded = 0
                    chunk_size = 1024 * 1024  # 1 MB chunks
                    
                    async for chunk in response.content.iter_chunked(chunk_size):
                        f.write(chunk)
                        downloaded += len(chunk)
                        
                        # Show progress for large files
                        if file_size and int(file_size) > 5 * 1024 * 1024:
                            progress = (downloaded / int(file_size)) * 100
                            print(f"Progress: {progress:.1f}%", end='\r')
                
                print()  # New line after progress
                
                # STEP 5: Verify file
                if output_path.exists():
                    actual_size = output_path.stat().st_size / (1024 * 1024)
                    # Return with download link for frontend
                    download_link = f"/api/files/download/{filename}"
                    return f"[OK] Successfully downloaded PDF!\n[FILE]: {filename}\n[DOWNLOAD_LINK]: {download_link}\n[SIZE]: {actual_size:.2f} MB"
                else:
                    return "[ERROR] Error: File was not saved properly"
    
    except aiohttp.ClientError as e:
        return f"[ERROR] Network error: {str(e)}"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error downloading PDF: {str(e)}"


@function_tool
async def batch_download_pdfs(urls_and_names: str) -> str:
    """
    Download multiple PDFs in PARALLEL for maximum speed.
    All downloads happen simultaneously, making it 5-10x faster than sequential downloads.
    
    Args:
        urls_and_names: JSON string of list of objects with 'url' and 'filename' keys.
                       Example: [{"url": "http://example.com/paper.pdf", "filename": "paper1"}]
    
    Returns:
        Summary of all download results
    """
    import json
    
    # Headers for bot detection bypass
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": "application/pdf,*/*",
        "Accept-Language": "en-US,en;q=0.9",
        "Referer": "https://scholar.google.com/",
    }
    downloads_folder = Path.home() / "Downloads"
    
    async def download_single_pdf(session, url, custom_filename):
        """Download a single PDF."""
        try:
            # Determine filename
            if custom_filename:
                filename = custom_filename if custom_filename.endswith('.pdf') else f"{custom_filename}.pdf"
            else:
                parsed_url = urlparse(url)
                filename = unquote(parsed_url.path.split('/')[-1])
                if not filename or not filename.endswith('.pdf'):
                    filename = "downloaded_document.pdf"
            
            # Clean filename
            filename = re.sub(r'[<>:"|?*]', '_', filename)
            output_path = downloads_folder / filename
            
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=120), ssl=False, allow_redirects=True) as response:
                if response.status != 200:
                    return f"[ERROR] HTTP {response.status}"
                
                with open(output_path, 'wb') as f:
                    async for chunk in response.content.iter_chunked(1024 * 512):
                        f.write(chunk)
                
                if output_path.exists():
                    return f"[OK] {filename}"
                return "[ERROR] Save failed"
        except Exception as e:
            return f"[ERROR] {str(e)[:50]}"
    
    try:
        pdf_list = json.loads(urls_and_names)
        if not isinstance(pdf_list, list) or len(pdf_list) == 0:
            return "[ERROR] Error: Invalid input"
        
        print(f"🚀 PARALLEL download starting: {len(pdf_list)} PDFs...")
        
        async with aiohttp.ClientSession(headers=headers) as session:
            tasks = [download_single_pdf(session, item.get('url', ''), item.get('filename', '')) for item in pdf_list]
            results = await asyncio.gather(*tasks, return_exceptions=True)
        
        success = sum(1 for r in results if isinstance(r, str) and "[OK]" in r)
        summary = [f"📦 BATCH DOWNLOAD: {success}/{len(pdf_list)} successful"]
        for i, result in enumerate(results):
            name = pdf_list[i].get('filename', f'pdf_{i+1}')
            summary.append(f"  {result if isinstance(result, str) else '[ERROR] Error'}")
        
        print(f"[OK] Batch complete: {success}/{len(pdf_list)}")
        return "\n".join(summary)
    
    except json.JSONDecodeError:
        return "[ERROR] Invalid JSON"
    except Exception as e:
        return f"[ERROR] Batch error: {str(e)}"

@function_tool
def extract_text_from_audio(
    audio_file_path: str,
    language: str = "en"
) -> str:
    """
    Extract text from audio files using Groq Whisper AI.
    Supports MP3, WAV, M4A, FLAC, OGG, and more formats.
    
    Args:
        audio_file_path: Path to the audio file
        language: Language code (default: "en" for English, "ur" for Urdu)
        
    Returns:
        Extracted text from the audio file
    """
    try:
        from groq import Groq
        from pathlib import Path
        
        audio_path = Path(audio_file_path)
        
        # Check if file exists
        if not audio_path.exists():
            return f"[ERROR] Error: Audio file not found at {audio_file_path}"
        
        # Supported formats by Whisper
        supported_formats = ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.webm', '.mp4', '.mpeg', '.mpga']
        if audio_path.suffix.lower() not in supported_formats:
            return f"[ERROR] Error: Unsupported format {audio_path.suffix}. Supported: {', '.join(supported_formats)}"
        
        # Get Groq API key
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "[ERROR] Error: groq_api_key not found in environment variables"
        
        # Initialize Groq client
        client = Groq(api_key=groq_key)
        
        print(f"🎤 Processing audio with Groq Whisper: {audio_path.name}")
        
        # Get file size
        file_size_mb = audio_path.stat().st_size / (1024 * 1024)
        print(f"📦 Audio Size: {file_size_mb:.2f} MB")
        
        # Check file size limit (25MB for Groq Whisper)
        if file_size_mb > 25:
            return f"[ERROR] Error: File too large ({file_size_mb:.2f} MB). Maximum size is 25 MB."
        
        # Open and transcribe audio file
        with open(audio_file_path, "rb") as audio_file:
            transcription = client.audio.transcriptions.create(
                model="whisper-large-v3-turbo",
                file=audio_file,
                language=language,
                response_format="text"
            )
        
        # Build result
        result = []
        result.append(f"🎤 Audio File: {audio_path.name}")
        result.append(f"📦 Size: {file_size_mb:.2f} MB")
        result.append(f"🌐 Language: {language}")
        result.append(f"🤖 Model: Groq Whisper Large V3 Turbo")
        result.append("=" * 80)
        result.append("")
        result.append("📝 Transcribed Text:")
        result.append(transcription)
        
        print(f"[OK] Successfully transcribed audio")
        return "\n".join(result)
    
    except ImportError:
        return "[ERROR] Error: groq is required. Install it with: pip install groq"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error processing audio file: {str(e)}"


@function_tool
def voice_output(
    text: str,
    save_file: bool = True,
    filename: str = "output.wav",
    voice_type: str = "natural"
) -> str:
    """
    Convert text to speech using Google Gemini Native Audio Dialog model.

    Args:
        text: The text to convert to speech
        save_file: Whether to save audio file (default: True)
        filename: Output filename (will be saved as .wav)
        voice_type: Voice style - "natural", "professional", "calm" (default: "natural")

    Returns:
        Success message with download link
    """
    try:
        from google import genai
        from google.genai import types
        import wave

        # Ensure filename ends with .wav
        if not filename.endswith('.wav'):
            filename = filename.rsplit('.', 1)[0] + '.wav'

        # Clean filename
        filename = re.sub(r'[<>:"|?*]', '_', filename)

        # Use project downloads folder (for web access)
        downloads_folder = Path(__file__).parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        output_path = downloads_folder / filename

        print(f"[AUDIO] Generating speech with Gemini Native Audio model...")

        # Get all API keys for rotation
        api_keys = [
            os.getenv("GEMINI_API_KEY_5"),
            os.getenv("GEMINI_API_KEY_4"),
            os.getenv("GEMINI_API_KEY_3"),
            os.getenv("GEMINI_API_KEY_2"),
            os.getenv("GEMINI_API_KEY_1"),
        ]
        api_keys = [k for k in api_keys if k]

        if not api_keys:
            return "[ERROR] No GEMINI_API_KEY found in environment variables"

        # Generate audio with API key rotation
        print(f"[AUDIO] Converting text to speech ({len(text)} characters)...")

        response = None
        last_error = None
        for key_idx, gemini_key in enumerate(api_keys):
            client = genai.Client(api_key=gemini_key)

            for retry in range(2):
                try:
                    print(f"[AUDIO] Trying API key {key_idx+1}/{len(api_keys)}...")
                    response = client.models.generate_content(
                        model="gemini-2.5-flash-preview-tts",
                        contents=text,
                        config=types.GenerateContentConfig(
                            response_modalities=["AUDIO"],
                            speech_config=types.SpeechConfig(
                                voice_config=types.VoiceConfig(
                                    prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                        voice_name='Aoede'
                                    )
                                )
                            ),
                        )
                    )
                    if response and response.candidates:
                        break
                except Exception as e:
                    last_error = str(e)
                    print(f"[AUDIO] API key {key_idx+1} error: {str(e)[:50]}...")
                    if "429" in str(e):
                        import time
                        time.sleep(3)
                    continue

            if response and response.candidates:
                break

        if not response or not response.candidates:
            return f"[ERROR] All API keys exhausted: {last_error}"

        # Get audio data (PCM format)
        if response.candidates and response.candidates[0].content.parts:
            audio_part = response.candidates[0].content.parts[0]
            if hasattr(audio_part, 'inline_data') and audio_part.inline_data:
                pcm_data = audio_part.inline_data.data

                # Save as WAV file with proper format
                with wave.open(str(output_path), "wb") as wf:
                    wf.setnchannels(1)  # Mono
                    wf.setsampwidth(2)  # 16-bit
                    wf.setframerate(24000)  # 24kHz sample rate
                    wf.writeframes(pcm_data)

                if output_path.exists():
                    file_size_kb = output_path.stat().st_size / 1024
                    download_link = f"/api/files/download/{filename}"
                    print(f"[OK] Audio file created: {file_size_kb:.2f} KB")
                    return f"[OK] Successfully created audio file!\n[FILE]: {filename}\n[DOWNLOAD_LINK]: {download_link}\n[SIZE]: {file_size_kb:.2f} KB"
                else:
                    return "[ERROR] Failed to save audio file"
            else:
                return "[ERROR] No audio data in response"
        else:
            return "[ERROR] No audio generated by model"

    except ImportError as ie:
        return f"[ERROR] Required package missing: {str(ie)}. Install with: pip install google-genai"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error creating audio: {str(e)}"


@function_tool
def read_folder(folder_path: str) -> str:
    """
    Read all PDF and Word files in a folder.
    
    Args:
        folder_path: The path to the folder containing files
        
    Returns:
        Content from all PDF and Word files in the folder
    """
    try:
        folder = Path(folder_path)
        
        if not folder.exists():
            return f"Error: Folder not found at {folder_path}"
        
        if not folder.is_dir():
            return f"Error: {folder_path} is not a folder"
        
        # Find all PDF and Word files
        pdf_files = list(folder.glob("*.pdf"))
        docx_files = list(folder.glob("*.docx"))
        
        all_files = pdf_files + docx_files
        
        if not all_files:
            return f"No PDF or Word files found in {folder_path}"
        
        results = []
        results.append(f"Found {len(pdf_files)} PDF files and {len(docx_files)} Word files\n")
        results.append("=" * 80)
        
        # Read each PDF file
        for pdf_file in pdf_files:
            results.append(f"\n\n📄 FILE: {pdf_file.name}")
            results.append("-" * 80)
            try:
                with open(pdf_file, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    num_pages = len(pdf_reader.pages)
                    
                    for page_num in range(num_pages):
                        page = pdf_reader.pages[page_num]
                        results.append(f"\nPage {page_num + 1}:\n{page.extract_text()}")
            except Exception as e:
                results.append(f"Error reading {pdf_file.name}: {str(e)}")
        
        # Read each Word file
        for docx_file in docx_files:
            results.append(f"\n\n📝 FILE: {docx_file.name}")
            results.append("-" * 80)
            try:
                doc = Document(docx_file)
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        results.append(paragraph.text)
                
                # Include tables
                for table_num, table in enumerate(doc.tables):
                    results.append(f"\nTable {table_num + 1}:")
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        results.append(row_text)
            except Exception as e:
                results.append(f"Error reading {docx_file.name}: {str(e)}")
        
        return "\n".join(results)
    
    except Exception as e:
        return f"Error reading folder: {str(e)}"


@function_tool
def list_files_in_folder(folder_path: str) -> str:
    """
    List all PDF and Word files in a folder without reading their content.
    
    Args:
        folder_path: The path to the folder
        
    Returns:
        List of file names and their sizes
    """
    try:
        folder = Path(folder_path)
        
        if not folder.exists():
            return f"Error: Folder not found at {folder_path}"
        
        if not folder.is_dir():
            return f"Error: {folder_path} is not a folder"
        
        # Find all PDF and Word files
        pdf_files = list(folder.glob("*.pdf"))
        docx_files = list(folder.glob("*.docx"))
        
        all_files = pdf_files + docx_files
        
        if not all_files:
            return f"No PDF or Word files found in {folder_path}"
        
        results = []
        results.append(f"📁 Folder: {folder_path}")
        results.append(f"Total files: {len(all_files)} ({len(pdf_files)} PDF, {len(docx_files)} Word)\n")
        
        for i, file in enumerate(all_files, 1):
            file_size = file.stat().st_size / 1024  # Size in KB
            file_type = "PDF" if file.suffix == ".pdf" else "Word"
            results.append(f"{i}. {file.name} ({file_type}, {file_size:.1f} KB)")
        
        return "\n".join(results)
    
    except Exception as e:
        return f"Error listing files: {str(e)}"


@function_tool
def create_word_file(content: str, file_name: str, title: str = "Document") -> str:
    """
    Create a Word document with formatted content including headings.
    
    Args:
        content: The text content to write in the document
        file_name: Name of the file to create (e.g., 'output.docx')
        title: Title/heading for the document (optional)
        
    Returns:
        Success message with file path
    """
    try:
        # Create a new Document
        doc = Document()
        
        # Add title
        title_paragraph = doc.add_heading(title, level=0)
        title_paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        # Split content by double newlines to create paragraphs
        paragraphs = content.split('\n\n')
        
        for para_text in paragraphs:
            para_text = para_text.strip()
            if not para_text:
                continue
            
            # Check if it looks like a heading (starts with #, all caps, or ends with :)
            if para_text.startswith('#'):
                # Markdown-style heading
                heading_level = para_text.count('#', 0, 3)
                heading_text = para_text.lstrip('#').strip()
                doc.add_heading(heading_text, level=min(heading_level, 3))
            elif para_text.isupper() and len(para_text.split()) <= 10:
                # All caps = heading
                doc.add_heading(para_text, level=1)
            elif para_text.endswith(':') and len(para_text.split()) <= 10 and '\n' not in para_text:
                # Ends with colon = subheading
                doc.add_heading(para_text.rstrip(':'), level=2)
            else:
                # Regular paragraph
                paragraph = doc.add_paragraph(para_text)
                paragraph.style = 'Normal'
        
        # Ensure file has .docx extension
        if not file_name.endswith('.docx'):
            file_name += '.docx'

        # Clean filename
        file_name = re.sub(r'[<>:"|?*]', '_', file_name)

        # Save to project downloads folder (for web access)
        downloads_folder = Path(__file__).parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        downloads_path = downloads_folder / file_name
        doc.save(str(downloads_path))

        file_size_kb = downloads_path.stat().st_size / 1024
        download_link = f"/api/files/download/{file_name}"
        return f"[OK] Successfully created Word document!\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: {download_link}\n[SIZE]: {file_size_kb:.2f} KB"

    except Exception as e:
        return f"[ERROR] Error creating Word document: {str(e)}"


@function_tool
def create_pdf(content: str, file_name: str, title: str = "Document") -> str:
    """
    Create a PDF document with formatted content including headings.
    
    Args:
        content: The text content to write in the PDF
        file_name: Name of the file to create (e.g., 'output.pdf')
        title: Title/heading for the document (optional)
        
    Returns:
        Success message with file path
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        
        # Ensure file has .pdf extension
        if not file_name.lower().endswith('.pdf'):
            file_name += '.pdf'

        # Clean filename
        file_name = re.sub(r'[<>:"|?*]', '_', file_name)

        # Save to project downloads folder (for web access)
        downloads_folder = Path(__file__).parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        downloads_path = downloads_folder / file_name
        
        # Create PDF document
        doc = SimpleDocTemplate(
            str(downloads_path),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        
        # Create custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            alignment=TA_CENTER,
            spaceAfter=30,
            textColor='#1a1a2e'
        )
        
        heading1_style = ParagraphStyle(
            'CustomHeading1',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=12,
            spaceBefore=20,
            textColor='#16213e'
        )
        
        heading2_style = ParagraphStyle(
            'CustomHeading2',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=10,
            spaceBefore=15,
            textColor='#0f3460'
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=8,
            alignment=TA_LEFT,
            leading=16
        )
        
        # Build content
        story = []
        
        # Add title
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3 * inch))
        
        # Split content by double newlines to create paragraphs
        paragraphs = content.split('\n\n')
        
        for para_text in paragraphs:
            para_text = para_text.strip()
            if not para_text:
                continue
            
            # Escape HTML special characters for ReportLab
            para_text = para_text.replace('&', '&amp;')
            para_text = para_text.replace('<', '&lt;')
            para_text = para_text.replace('>', '&gt;')
            
            # Check if it looks like a heading
            if para_text.startswith('#'):
                # Markdown-style heading
                heading_level = para_text.count('#', 0, 3)
                heading_text = para_text.lstrip('#').strip()
                if heading_level == 1:
                    story.append(Paragraph(heading_text, heading1_style))
                else:
                    story.append(Paragraph(heading_text, heading2_style))
            elif para_text.isupper() and len(para_text.split()) <= 10:
                # All caps = heading
                story.append(Paragraph(para_text, heading1_style))
            elif para_text.endswith(':') and len(para_text.split()) <= 10 and '\n' not in para_text:
                # Ends with colon = subheading
                story.append(Paragraph(para_text.rstrip(':'), heading2_style))
            else:
                # Regular paragraph - handle line breaks
                para_text = para_text.replace('\n', '<br/>')
                story.append(Paragraph(para_text, body_style))
        
        # Build PDF
        doc.build(story)

        file_size_kb = downloads_path.stat().st_size / 1024
        download_link = f"/api/files/download/{file_name}"
        return f"[OK] Successfully created PDF document!\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: {download_link}\n[SIZE]: {file_size_kb:.2f} KB"
    
    except ImportError:
        return "[ERROR] Error: reportlab is required. Install it with: pip install reportlab"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error creating PDF document: {str(e)}"


@function_tool
def create_pptx(content: str, file_name: str = "presentation.pptx", theme: str = "professional") -> str:
    """
    Create a professional PowerPoint presentation from text content.
    Uses AI to structure content into slides with themes and designs.

    Args:
        content: The text content to convert into a presentation
        file_name: Name of the file to create (e.g., 'output.pptx')
        theme: Theme style - "professional" (blue), "modern" (dark), "elegant" (purple), "nature" (green), "warm" (orange)

    Returns:
        Success message with download link
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import json
        from google import genai

        # Theme color schemes
        themes = {
            "professional": {
                "bg_color": RGBColor(0x1a, 0x1a, 0x2e),      # Dark blue
                "title_color": RGBColor(0xff, 0xff, 0xff),   # White
                "text_color": RGBColor(0xe0, 0xe0, 0xe0),    # Light gray
                "accent_color": RGBColor(0x00, 0xd4, 0xff),  # Cyan
                "subtitle_color": RGBColor(0x90, 0xca, 0xf9) # Light blue
            },
            "modern": {
                "bg_color": RGBColor(0x12, 0x12, 0x12),      # Dark
                "title_color": RGBColor(0xff, 0xff, 0xff),   # White
                "text_color": RGBColor(0xcc, 0xcc, 0xcc),    # Gray
                "accent_color": RGBColor(0xff, 0x61, 0x61),  # Red accent
                "subtitle_color": RGBColor(0xff, 0x99, 0x99) # Light red
            },
            "elegant": {
                "bg_color": RGBColor(0x2d, 0x1b, 0x4e),      # Dark purple
                "title_color": RGBColor(0xff, 0xff, 0xff),   # White
                "text_color": RGBColor(0xe0, 0xd0, 0xf0),    # Light lavender
                "accent_color": RGBColor(0xbb, 0x86, 0xfc),  # Purple accent
                "subtitle_color": RGBColor(0xce, 0x93, 0xd8) # Light purple
            },
            "nature": {
                "bg_color": RGBColor(0x1b, 0x2e, 0x1b),      # Dark green
                "title_color": RGBColor(0xff, 0xff, 0xff),   # White
                "text_color": RGBColor(0xd0, 0xe8, 0xd0),    # Light green
                "accent_color": RGBColor(0x4c, 0xaf, 0x50),  # Green accent
                "subtitle_color": RGBColor(0xa5, 0xd6, 0xa7) # Light green
            },
            "warm": {
                "bg_color": RGBColor(0x2e, 0x1a, 0x0a),      # Dark orange/brown
                "title_color": RGBColor(0xff, 0xff, 0xff),   # White
                "text_color": RGBColor(0xf5, 0xe0, 0xd0),    # Light peach
                "accent_color": RGBColor(0xff, 0x98, 0x00),  # Orange accent
                "subtitle_color": RGBColor(0xff, 0xcc, 0x80) # Light orange
            }
        }

        # Get theme colors
        theme_colors = themes.get(theme.lower(), themes["professional"])

        print(f"[PPTX] Creating PowerPoint presentation with {theme} theme...")

        # Use Gemini to structure content into slides
        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found for content structuring"

        client = genai.Client(api_key=gemini_key)

        # Prompt Gemini to structure content
        structure_prompt = f"""Analyze this text and structure it into PowerPoint slides.
Return a JSON object with this exact format:
{{
    "title": "Main Presentation Title",
    "subtitle": "Optional subtitle",
    "slides": [
        {{
            "title": "Slide Title",
            "bullet_points": ["Point 1", "Point 2", "Point 3"]
        }}
    ]
}}

Rules:
- Create 3-8 slides based on content length
- Each slide should have 2-5 bullet points
- Keep bullet points concise (under 15 words each)
- First slide should be a title slide (no bullet points needed)
- Last slide can be a summary/conclusion

TEXT TO STRUCTURE:
{content}

Return ONLY the JSON object, no other text."""

        print("[PPTX] Using AI to structure content into slides...")

        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=structure_prompt
            )

            response_text = response.text.strip()
            # Clean up response - remove markdown code blocks if present
            if response_text.startswith("```"):
                response_text = response_text.split("```")[1]
                if response_text.startswith("json"):
                    response_text = response_text[4:]
            response_text = response_text.strip()

            slide_data = json.loads(response_text)
        except Exception as e:
            print(f"[PPTX] AI structuring failed: {e}, using fallback...")
            # Fallback: Create simple slides from paragraphs
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            slide_data = {
                "title": "Presentation",
                "subtitle": "",
                "slides": []
            }

            # Group paragraphs into slides (2-3 per slide)
            for i in range(0, len(paragraphs), 3):
                chunk = paragraphs[i:i+3]
                slide_data["slides"].append({
                    "title": f"Section {len(slide_data['slides']) + 1}",
                    "bullet_points": chunk
                })

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)

        def set_slide_background(slide, color):
            """Set solid background color for a slide."""
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_text_box(slide, left, top, width, height, text, font_size, font_color, bold=False, alignment=PP_ALIGN.LEFT):
            """Add a text box to the slide."""
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.alignment = alignment
            return txBox

        # Create title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        set_slide_background(title_slide, theme_colors["bg_color"])

        # Add title
        add_text_box(
            title_slide,
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5),
            slide_data.get("title", "Presentation"),
            48, theme_colors["title_color"], bold=True, alignment=PP_ALIGN.CENTER
        )

        # Add subtitle if present
        if slide_data.get("subtitle"):
            add_text_box(
                title_slide,
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8),
                slide_data["subtitle"],
                24, theme_colors["subtitle_color"], alignment=PP_ALIGN.CENTER
            )

        # Add accent line
        line = title_slide.shapes.add_shape(
            1,  # Rectangle
            Inches(4), Inches(4), Inches(5.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme_colors["accent_color"]
        line.line.fill.background()

        # Create content slides
        for slide_info in slide_data.get("slides", []):
            content_slide = prs.slides.add_slide(title_slide_layout)
            set_slide_background(content_slide, theme_colors["bg_color"])

            # Add slide title
            add_text_box(
                content_slide,
                Inches(0.5), Inches(0.4), Inches(12.333), Inches(1),
                slide_info.get("title", ""),
                36, theme_colors["title_color"], bold=True
            )

            # Add accent line under title
            accent_line = content_slide.shapes.add_shape(
                1, Inches(0.5), Inches(1.3), Inches(2), Inches(0.04)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = theme_colors["accent_color"]
            accent_line.line.fill.background()

            # Add bullet points
            bullet_points = slide_info.get("bullet_points", [])
            y_position = 1.8

            for point in bullet_points:
                if not point.strip():
                    continue

                # Add bullet marker
                bullet_box = content_slide.shapes.add_textbox(
                    Inches(0.5), Inches(y_position), Inches(0.3), Inches(0.4)
                )
                bullet_tf = bullet_box.text_frame
                bullet_p = bullet_tf.paragraphs[0]
                bullet_p.text = "●"
                bullet_p.font.size = Pt(14)
                bullet_p.font.color.rgb = theme_colors["accent_color"]

                # Add point text
                add_text_box(
                    content_slide,
                    Inches(0.9), Inches(y_position), Inches(11.5), Inches(0.8),
                    point,
                    20, theme_colors["text_color"]
                )

                y_position += 0.9

        # Ensure file has .pptx extension
        if not file_name.lower().endswith('.pptx'):
            file_name += '.pptx'

        # Clean filename
        file_name = re.sub(r'[<>:"|?*]', '_', file_name)

        # Save to downloads folder
        downloads_folder = Path(__file__).parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        output_path = downloads_folder / file_name

        prs.save(str(output_path))

        if output_path.exists():
            file_size_kb = output_path.stat().st_size / 1024
            download_link = f"/api/files/download/{file_name}"
            num_slides = len(prs.slides)
            print(f"[OK] Created PowerPoint: {num_slides} slides, {file_size_kb:.2f} KB")
            return f"[OK] Successfully created PowerPoint presentation!\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: {download_link}\n[SIZE]: {file_size_kb:.2f} KB\n[SLIDES]: {num_slides}\n[THEME]: {theme}"
        else:
            return "[ERROR] Failed to save PowerPoint file"

    except ImportError as ie:
        missing = str(ie)
        if "pptx" in missing.lower():
            return "[ERROR] python-pptx required. Install: pip install python-pptx"
        elif "lxml" in missing.lower():
            return "[ERROR] lxml required. Install: pip install lxml"
        else:
            return f"[ERROR] Missing library: {missing}"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error creating PowerPoint: {str(e)}"


@function_tool
def read_pdf(file_path: str) -> str:
    """
    Read and extract text from a PDF file using PyPDF2.

    Args:
        file_path: The path to the PDF file (can be just filename if in downloads folder)

    Returns:
        Extracted text content from the PDF with page numbers clearly marked
    """
    try:
        # If just filename, check downloads folder first
        pdf_path = Path(file_path)
        if not pdf_path.exists() and not pdf_path.is_absolute():
            downloads_folder = Path(__file__).parent / "downloads"
            pdf_path = downloads_folder / file_path
            if not pdf_path.suffix:
                pdf_path = pdf_path.with_suffix('.pdf')

        file_path = str(pdf_path)

        # Check if file exists
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"

        print(f"[READ] Reading PDF: {Path(file_path).name}")

        # Use PyPDF2 for text extraction
        import PyPDF2

        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            num_pages = len(pdf_reader.pages)

            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            print(f"[INFO] PDF Size: {file_size_mb:.2f} MB, Pages: {num_pages}")

            text_content = []
            text_content.append(f"PDF Document: {Path(file_path).name}")
            text_content.append(f"Total Pages: {num_pages}")
            text_content.append("=" * 60)

            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text() or ""

                text_content.append(f"\nPAGE {page_num + 1}:\n")
                text_content.append(page_text if page_text.strip() else "[No text on this page]")
                text_content.append(f"\n[End of Page {page_num + 1}]")

            result = "\n".join(text_content)
            print(f"[OK] Extracted {len(result)} characters from PDF")
            return result

    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error reading PDF: {str(e)}"


@function_tool
def read_image(file_path: str) -> str:
    """
    Read and extract text from an image using Groq Llama 4 Scout Vision OCR.
    Supports PNG, JPG, JPEG, GIF, BMP, WEBP formats.
    
    Args:
        file_path: The path to the image file to read
        
    Returns:
        Extracted text content from the image
    """
    try:
        from groq import Groq
        import base64
        
        # Check if file exists
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        # Get Groq API key
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "[ERROR] Error: groq_api_key not found in environment variables"
        
        # Initialize Groq client
        client = Groq(api_key=groq_key)
        
        # Get file extension
        file_ext = Path(file_path).suffix.lower()
        supported_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']
        
        if file_ext not in supported_formats:
            return f"[ERROR] Error: Unsupported format {file_ext}. Supported: {', '.join(supported_formats)}"
        
        # Determine mime type
        mime_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp'
        }
        mime_type = mime_types.get(file_ext, 'image/png')
        
        print(f"🖼️ Processing image with Groq Llama 4 Scout OCR: {Path(file_path).name}")
        
        # Read and encode image
        with open(file_path, 'rb') as f:
            img_data = f.read()
        
        img_base64 = base64.b64encode(img_data).decode()
        file_size_kb = len(img_data) / 1024
        
        print(f"📦 Image Size: {file_size_kb:.2f} KB")
        
        # Send to Groq Llama 4 Scout (vision model)
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract ALL text from this image. Perform accurate OCR. Preserve formatting and structure. Return ONLY the extracted text, nothing else."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{img_base64}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4096
        )
        
        extracted_text = response.choices[0].message.content if response.choices else "[No text extracted]"
        
        # Build output
        text_content = []
        text_content.append(f"🖼️ Image: {Path(file_path).name}")
        text_content.append(f"📦 Size: {file_size_kb:.2f} KB")
        text_content.append(f"🤖 OCR Method: Groq Llama 4 Scout Vision")
        text_content.append("=" * 80)
        text_content.append("")
        text_content.append(extracted_text)
        
        print(f"[OK] Successfully extracted text from image")
        return "\n".join(text_content)
    
    except ImportError as ie:
        missing_lib = str(ie)
        if "groq" in missing_lib.lower():
            return "[ERROR] Error: groq is required. Install it with: pip install groq"
        else:
            return f"[ERROR] Error: Missing library - {missing_lib}"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error reading image with OCR: {str(e)}"


@function_tool
def read_pptx(file_path: str) -> str:
    """
    Read and extract text from a PowerPoint presentation using Groq Llama 4 Scout Vision OCR.
    Uses cloud AI to read slides - no local CPU/GPU required.
    
    Args:
        file_path: The path to the PowerPoint file (.pptx)
        
    Returns:
        Extracted text content from all slides with slide numbers
    """
    try:
        from groq import Groq
        from pptx import Presentation
        from pptx.util import Inches
        import base64
        from io import BytesIO
        
        # Check if file exists
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        # Check file extension
        if not file_path.lower().endswith(('.pptx', '.ppt')):
            return f"[ERROR] Error: File must be a PowerPoint file (.pptx or .ppt)"
        
        # Get Groq API key
        groq_key = os.getenv("groq_api_key", "")
        if not groq_key:
            return "[ERROR] Error: groq_api_key not found in environment variables"
        
        # Initialize Groq client
        client = Groq(api_key=groq_key)
        
        print(f"📊 Processing PowerPoint with Groq AI: {Path(file_path).name}")
        
        # Load presentation
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        print(f"📦 File Size: {file_size_mb:.2f} MB, Slides: {num_slides}")
        
        # Extract text from all slides
        text_content = []
        text_content.append(f"📊 PowerPoint: {Path(file_path).name}")
        text_content.append(f"📦 Size: {file_size_mb:.2f} MB")
        text_content.append(f"📑 Total Slides: {num_slides}")
        text_content.append(f"🤖 Method: Text Extraction + Groq AI for complex content")
        text_content.append("=" * 80)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"📖 Processing slide {slide_num} of {num_slides}...")
            
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Handle tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(f"[Table Row] {row_text}")
            
            # Add slide header
            text_content.append(f"\n\n{'#'*80}")
            text_content.append(f"#{'':^78}#")
            text_content.append(f"#{'📑 SLIDE ' + str(slide_num) + ' OF ' + str(num_slides):^78}#")
            text_content.append(f"#{'':^78}#")
            text_content.append(f"{'#'*80}\n")
            
            if slide_text:
                text_content.append("\n".join(slide_text))
            else:
                text_content.append("[No text content on this slide - may contain images/graphics]")
            
            text_content.append(f"\n{'─'*80}")
            text_content.append(f"[End of Slide {slide_num}]")
            text_content.append(f"{'─'*80}")
        
        print(f"[OK] Successfully processed all {num_slides} slides")
        return "\n".join(text_content)
    
    except ImportError as ie:
        missing_lib = str(ie)
        if "pptx" in missing_lib.lower():
            return "[ERROR] Error: python-pptx is required. Install it with: pip install python-pptx"
        elif "groq" in missing_lib.lower():
            return "[ERROR] Error: groq is required. Install it with: pip install groq"
        else:
            return f"[ERROR] Error: Missing library - {missing_lib}"
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error reading PowerPoint: {str(e)}"


@function_tool
def read_word(file_path: str) -> str:
    """
    Read and extract text from a Word document (.docx file) with page tracking.
    
    Args:
        file_path: The path to the Word document to read
        
    Returns:
        Extracted text content from the Word document with page information
    """
    try:
        # Check if file exists
        if not Path(file_path).exists():
            return f"Error: File not found at {file_path}"
        
        # Check if it's a .docx file
        if not file_path.lower().endswith('.docx'):
            return f"Error: File must be a .docx file. Old .doc format is not supported."
        
        # Open and read the Word document
        doc = Document(file_path)
        
        # Build structured output
        output = []
        output.append(f"📝 Word Document: {Path(file_path).name}")
        output.append("=" * 80)
        
        # Track current page (approximation based on page breaks)
        current_page = 1
        paragraphs_on_page = []
        
        for i, paragraph in enumerate(doc.paragraphs):
            # Check if paragraph contains a page break
            if paragraph._element.xpath('.//w:br[@w:type="page"]'):
                # Save current page content
                if paragraphs_on_page:
                    output.append(f"\n{'='*80}")
                    output.append(f"📖 PAGE {current_page}")
                    output.append(f"{'='*80}\n")
                    output.append("\n".join(paragraphs_on_page))
                    paragraphs_on_page = []
                
                current_page += 1
            
            # Add paragraph text if not empty
            if paragraph.text.strip():
                paragraphs_on_page.append(paragraph.text)
        
        # Add last page content
        if paragraphs_on_page:
            output.append(f"\n{'='*80}")
            output.append(f"📖 PAGE {current_page}")
            output.append(f"{'='*80}\n")
            output.append("\n".join(paragraphs_on_page))
        
        # Extract text from tables
        if doc.tables:
            output.append(f"\n{'='*80}")
            output.append(f"📊 TABLES ({len(doc.tables)} found)")
            output.append(f"{'='*80}\n")
            
            for table_num, table in enumerate(doc.tables, 1):
                output.append(f"\n--- Table {table_num} ---")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    output.append(row_text)
        
        # Summary
        num_paragraphs = sum(1 for p in doc.paragraphs if p.text.strip())
        summary = f"\n\n📊 Document Summary:\n"
        summary += f"- Approximate Pages: {current_page}\n"
        summary += f"- Total Paragraphs: {num_paragraphs}\n"
        summary += f"- Total Tables: {len(doc.tables)}"
        
        output.insert(1, summary)
        
        return "\n".join(output)
    
    except Exception as e:
        return f"Error reading Word document: {str(e)}"

@function_tool
def delete_pdf(file_path: str) -> str:
    """
    Delete a PDF file from the project downloads folder.

    Args:
        file_path: Filename of the PDF to delete (e.g., "paper.pdf")

    Returns:
        Success message or error
    """
    try:
        from pathlib import Path
        import os

        # Convert to Path object
        file = Path(file_path)

        # If just filename provided, check project downloads folder
        if not file.is_absolute():
            downloads_folder = Path(__file__).parent / "downloads"
            file = downloads_folder / file_path

            # Add .pdf if not present
            if not file.suffix:
                file = file.with_suffix('.pdf')

        # Check if file exists
        if not file.exists():
            return f"[ERROR] File not found: {file.name}"

        # Get file info before deletion
        file_name = file.name

        # Delete the file
        os.remove(file)
        
        # Verify deletion
        if not file.exists():
            return f"[OK] Successfully deleted: {file_name}"
        else:
            return "[ERROR] Error: File still exists after deletion attempt"
    
    except PermissionError:
        return f"[ERROR] Permission denied: Cannot delete {file.name}\n\nFile may be open in another program. Please close it and try again."
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error details: {error_details}")
        return f"[ERROR] Error deleting file: {str(e)}"


# ============================================================================
# RESEARCH ASSISTANT TOOLS - Advanced Features for Researchers
# ============================================================================

@function_tool
def smart_summarize_paper(paper_content: str, summary_type: str = "comprehensive") -> str:
    """
    Generate intelligent summaries of research papers with multiple summary types.

    Args:
        paper_content: The full text content of the paper (from read_pdf)
        summary_type: Type of summary:
            - "comprehensive" - Full analysis with all sections
            - "abstract" - Brief 2-3 sentence summary
            - "key_points" - Bullet points of main contributions
            - "methodology" - Focus on research methods
            - "beginner" - Explain like I'm new to this field

    Returns:
        Structured summary based on the requested type
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        prompts = {
            "comprehensive": f"""Analyze this research paper and provide a comprehensive summary:

PAPER CONTENT:
{paper_content[:15000]}

Provide the following sections:

## 1. ABSTRACT SUMMARY (2-3 sentences)
Brief overview of what the paper is about.

## 2. KEY CONTRIBUTIONS (3-5 bullet points)
- Main findings and contributions

## 3. METHODOLOGY
- Research design
- Data collection methods
- Analysis techniques

## 4. MAIN FINDINGS
- Key results
- Statistical significance (if applicable)

## 5. LIMITATIONS
- Acknowledged limitations
- Potential weaknesses

## 6. IMPLICATIONS
- Practical implications
- Future research directions

## 7. CITATION INFO
Extract: Title, Authors, Year, Journal/Conference (if available)

Format the output clearly with markdown headers.""",

            "abstract": f"""Summarize this research paper in 2-3 concise sentences that capture the main objective, methodology, and key findings:

{paper_content[:10000]}

Return ONLY the summary, nothing else.""",

            "key_points": f"""Extract the KEY CONTRIBUTIONS and MAIN POINTS from this research paper as bullet points:

{paper_content[:12000]}

Format:
## Key Contributions
- [Point 1]
- [Point 2]
- [Point 3]
- [Point 4]
- [Point 5]

## Main Findings
- [Finding 1]
- [Finding 2]
- [Finding 3]

Be concise but informative.""",

            "methodology": f"""Analyze the METHODOLOGY section of this research paper in detail:

{paper_content[:12000]}

Provide:
## Research Design
- Type of study (experimental, observational, qualitative, etc.)

## Data Collection
- How was data collected?
- Sample size and characteristics

## Analysis Methods
- Statistical methods used
- Tools/software mentioned

## Validity & Reliability
- How did authors ensure validity?

## Limitations of Methodology
- What are the methodological weaknesses?""",

            "beginner": f"""Explain this research paper in simple terms for someone NEW to this field:

{paper_content[:12000]}

Use:
- Simple language (no jargon)
- Analogies where helpful
- Short sentences
- Explain technical terms when used

Structure:
## What is this paper about? (1-2 sentences)

## Why does this matter? (Real-world importance)

## What did the researchers do? (Methods in simple terms)

## What did they find? (Key results)

## What does this mean for the field?

## Key terms explained:
- [Term 1]: [Simple explanation]
- [Term 2]: [Simple explanation]"""
        }

        prompt = prompts.get(summary_type, prompts["comprehensive"])

        print(f"[SUMMARIZE] Generating {summary_type} summary...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"PAPER SUMMARY - Type: {summary_type.upper()}")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Summary generated successfully")
        return "\n".join(output)

    except Exception as e:
        import traceback
        print(f"Error: {traceback.format_exc()}")
        return f"[ERROR] Error generating summary: {str(e)}"


@function_tool
def generate_citation(paper_content: str, citation_style: str = "all") -> str:
    """
    Generate formatted citations from paper content in multiple styles.

    Args:
        paper_content: The paper content or metadata
        citation_style: Citation format:
            - "bibtex" - BibTeX format
            - "apa" - APA 7th edition
            - "mla" - MLA format
            - "harvard" - Harvard style
            - "chicago" - Chicago style
            - "ieee" - IEEE format
            - "all" - Generate all formats

    Returns:
        Formatted citation(s)
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        prompt = f"""Extract bibliographic information from this paper and generate citations:

PAPER CONTENT (first part):
{paper_content[:5000]}

First, extract:
- Title
- Authors (all of them)
- Year
- Journal/Conference name
- Volume, Issue, Pages (if available)
- DOI (if available)
- Publisher (if available)

Then generate citations in these formats:
"""

        if citation_style == "all":
            prompt += """
## BibTeX
```bibtex
@article{...}
```

## APA 7th Edition
[Full APA citation]

## MLA
[Full MLA citation]

## Harvard
[Full Harvard citation]

## Chicago
[Full Chicago citation]

## IEEE
[Full IEEE citation]
"""
        else:
            style_prompts = {
                "bibtex": "Generate ONLY a BibTeX citation:\n```bibtex\n@article{...}\n```",
                "apa": "Generate ONLY an APA 7th edition citation.",
                "mla": "Generate ONLY an MLA format citation.",
                "harvard": "Generate ONLY a Harvard style citation.",
                "chicago": "Generate ONLY a Chicago style citation.",
                "ieee": "Generate ONLY an IEEE format citation."
            }
            prompt += style_prompts.get(citation_style, style_prompts["apa"])

        print(f"[CITATION] Generating {citation_style} citation(s)...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"GENERATED CITATIONS")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Citations generated successfully")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error generating citation: {str(e)}"


@function_tool
def compare_papers(papers_content: str) -> str:
    """
    Compare multiple research papers and identify agreements, disagreements, and gaps.

    Args:
        papers_content: Combined content of multiple papers, separated by "---PAPER---"
                       Format: "Paper 1 content ---PAPER--- Paper 2 content ---PAPER--- Paper 3 content"

    Returns:
        Detailed comparison analysis
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        prompt = f"""Analyze and compare these research papers:

{papers_content[:20000]}

Provide a detailed comparison:

## 1. PAPER IDENTIFICATION
For each paper, identify:
- Title
- Authors
- Year
- Main focus

## 2. AGREEMENTS
What do these papers agree on?
- Common findings
- Shared conclusions
- Similar methodologies

## 3. DISAGREEMENTS
Where do these papers differ?
- Conflicting findings
- Different interpretations
- Methodological differences

## 4. METHODOLOGICAL COMPARISON
| Aspect | Paper 1 | Paper 2 | Paper 3 |
|--------|---------|---------|---------|
| Design | | | |
| Sample | | | |
| Methods| | | |

## 5. RESEARCH GAPS IDENTIFIED
Based on these papers, what gaps exist in the literature?
- Unexplored areas
- Contradictions needing resolution
- Future research opportunities

## 6. SYNTHESIS
How do these papers together contribute to understanding the topic?

## 7. RECOMMENDATIONS
Which paper(s) should be prioritized for:
- Understanding fundamentals
- Latest developments
- Methodological guidance"""

        print(f"[COMPARE] Analyzing multiple papers...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"MULTI-PAPER COMPARISON ANALYSIS")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Comparison analysis completed")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error comparing papers: {str(e)}"


@function_tool
def write_literature_review(papers_content: str, topic: str = "", style: str = "academic") -> str:
    """
    Generate a literature review section from multiple papers.

    Args:
        papers_content: Combined content of papers (separated by ---PAPER---)
        topic: The topic/theme of the literature review
        style: Writing style - "academic", "concise", "detailed"

    Returns:
        Formatted literature review with proper citations
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        style_instructions = {
            "academic": "Use formal academic language, passive voice where appropriate, and scholarly tone.",
            "concise": "Be brief and to the point. Focus on key findings only.",
            "detailed": "Provide comprehensive coverage with detailed explanations and connections."
        }

        topic_text = f"on the topic of '{topic}'" if topic else ""

        prompt = f"""Write a LITERATURE REVIEW section {topic_text} based on these research papers:

{papers_content[:20000]}

INSTRUCTIONS:
- {style_instructions.get(style, style_instructions["academic"])}
- Synthesize findings across papers (don't just summarize each paper separately)
- Use proper in-text citations (Author, Year) format
- Group related findings thematically
- Identify trends, patterns, and gaps
- Maintain logical flow between paragraphs

STRUCTURE:
## Literature Review

### Introduction (1 paragraph)
Brief overview of the research landscape.

### Theme 1: [Identify main theme from papers]
Synthesize findings from multiple papers on this theme.

### Theme 2: [Second theme]
Synthesize findings on this theme.

### Theme 3: [Third theme if applicable]
Synthesize findings on this theme.

### Research Gaps and Future Directions
Based on the reviewed literature, identify gaps.

### Summary
Brief synthesis of the main takeaways.

---
## References
List all papers cited in the review (APA format).

IMPORTANT:
- Never fabricate citations
- Only cite papers actually provided
- Use (Author, Year) format for in-text citations"""

        print(f"[LIT REVIEW] Writing literature review ({style} style)...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"GENERATED LITERATURE REVIEW")
        output.append(f"Topic: {topic if topic else 'Based on provided papers'}")
        output.append(f"Style: {style}")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Literature review generated")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error writing literature review: {str(e)}"


@function_tool
def refine_research_question(topic: str, context: str = "") -> str:
    """
    Help refine a vague research idea into clear research questions, hypotheses, and variables.

    Args:
        topic: The research topic or vague idea
        context: Additional context (field, constraints, available resources)

    Returns:
        Refined research questions with hypotheses and methodology suggestions
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        context_text = f"\nAdditional context: {context}" if context else ""

        prompt = f"""Help refine this research topic into clear research questions:

TOPIC: {topic}
{context_text}

Provide:

## 1. TOPIC ANALYSIS
- What is the core phenomenon being studied?
- Why is this important?
- What gaps might exist?

## 2. REFINED RESEARCH QUESTIONS

### Primary Research Question (RQ1)
[Clear, focused, answerable question]

### Secondary Research Questions
- RQ2: [Related question]
- RQ3: [Related question]

## 3. HYPOTHESES (if applicable)
- H1: [Testable hypothesis for RQ1]
- H2: [Alternative hypothesis]

## 4. VARIABLES

### For Quantitative Research:
| Variable Type | Variable Name | Measurement |
|---------------|---------------|-------------|
| Independent   |               |             |
| Dependent     |               |             |
| Control       |               |             |

### For Qualitative Research:
- Key concepts to explore
- Themes to investigate

## 5. SUGGESTED METHODOLOGY

### Approach Options:
1. **Quantitative**: [Specific method suggestion]
2. **Qualitative**: [Specific method suggestion]
3. **Mixed Methods**: [How to combine]

### Data Collection Suggestions:
- Primary data: [Methods]
- Secondary data: [Sources]

### Analysis Techniques:
- [Suggested analytical approaches]

## 6. POTENTIAL CHALLENGES
- [Challenge 1 and mitigation]
- [Challenge 2 and mitigation]

## 7. RECOMMENDED READING
Suggest 3-5 seminal papers/books to start with (describe what to search for)."""

        print(f"[RESEARCH Q] Refining research question...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"RESEARCH QUESTION REFINEMENT")
        output.append(f"Original Topic: {topic}")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Research question refined")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error refining research question: {str(e)}"


@function_tool
def extract_paper_metadata(paper_content: str) -> str:
    """
    Extract structured metadata from a research paper.

    Args:
        paper_content: The paper content

    Returns:
        Structured metadata (title, authors, abstract, keywords, etc.)
    """
    try:
        from google import genai
        import json

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        prompt = f"""Extract metadata from this research paper:

{paper_content[:8000]}

Return a structured extraction:

## METADATA

**Title:** [Full title]

**Authors:** [All authors, comma-separated]

**Year:** [Publication year]

**Journal/Conference:** [Where published]

**Volume/Issue/Pages:** [If available]

**DOI:** [If available]

**Abstract:** [Full abstract]

**Keywords:** [Listed keywords or extracted key terms]

**Research Type:** [Empirical/Theoretical/Review/Meta-analysis/etc.]

**Field/Discipline:** [Primary field]

**Methodology:** [Brief: Quantitative/Qualitative/Mixed/etc.]

**Sample/Data:** [Brief description of data used]

**Key Findings:** [1-2 sentences]

**Limitations Mentioned:** [Brief]

**Future Work Suggested:** [Brief]"""

        print(f"[METADATA] Extracting paper metadata...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"EXTRACTED PAPER METADATA")
        output.append(f"{'='*80}\n")
        output.append(result)

        print(f"[OK] Metadata extracted")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error extracting metadata: {str(e)}"


@function_tool
def write_section(content: str, section_type: str, style: str = "academic") -> str:
    """
    Help write specific sections of a research paper.

    Args:
        content: Source content/notes/data to use
        section_type: Which section to write:
            - "abstract" - Write abstract
            - "introduction" - Write introduction
            - "related_work" - Write related work/literature review
            - "methodology" - Write methodology section
            - "results" - Help structure results
            - "discussion" - Write discussion
            - "conclusion" - Write conclusion
        style: "academic", "concise", "journal" (specific journal style)

    Returns:
        Draft of the requested section
    """
    try:
        from google import genai

        gemini_key = os.getenv("GEMINI_API_KEY_5") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "[ERROR] No GEMINI_API_KEY found"

        client = genai.Client(api_key=gemini_key)

        section_prompts = {
            "abstract": f"""Write an ABSTRACT for a research paper based on this content:

{content[:10000]}

The abstract should:
- Be 150-300 words
- Include: Background, Objective, Methods, Results, Conclusion
- Be self-contained
- Use past tense for methods/results
- Avoid citations and abbreviations

Write in {style} style.""",

            "introduction": f"""Write an INTRODUCTION section based on this content:

{content[:12000]}

Structure:
1. Opening hook - Why is this topic important?
2. Background - What do we know?
3. Gap - What's missing in current knowledge?
4. Objective - What does this research aim to do?
5. Contribution - What's new/significant?
6. Paper structure (optional) - Brief roadmap

Write in {style} style with proper academic tone.""",

            "related_work": f"""Write a RELATED WORK / LITERATURE REVIEW section:

{content[:15000]}

Structure:
- Organize thematically (not paper by paper)
- Show evolution of research in this area
- Identify trends and patterns
- Highlight gaps your work addresses
- Use proper citations (Author, Year)

Write in {style} style.""",

            "methodology": f"""Write a METHODOLOGY section based on:

{content[:12000]}

Include:
1. Research Design - Type of study
2. Participants/Data - Sample description
3. Materials/Instruments - Tools used
4. Procedure - Step by step process
5. Data Analysis - How data was analyzed
6. Ethical Considerations (if applicable)

Be specific enough for replication. Write in {style} style.""",

            "results": f"""Help structure a RESULTS section based on:

{content[:12000]}

Organize:
1. Overview of findings
2. Present results logically (by research question or hypothesis)
3. Include statistical details where relevant
4. Reference tables/figures
5. Report effect sizes and confidence intervals

Write objectively without interpretation. Use {style} style.""",

            "discussion": f"""Write a DISCUSSION section based on:

{content[:12000]}

Structure:
1. Summary of key findings
2. Interpretation - What do results mean?
3. Comparison with previous research
4. Theoretical implications
5. Practical implications
6. Limitations
7. Future research directions

Write in {style} style.""",

            "conclusion": f"""Write a CONCLUSION section based on:

{content[:10000]}

Include:
1. Restate the research problem
2. Summarize main findings (without new info)
3. State the significance
4. Final thoughts / call to action

Keep it concise (1-2 paragraphs). Write in {style} style."""
        }

        prompt = section_prompts.get(section_type, section_prompts["abstract"])

        print(f"[WRITING] Drafting {section_type} section...")

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )

        result = response.text.strip()

        output = []
        output.append(f"{'='*80}")
        output.append(f"DRAFT: {section_type.upper()} SECTION")
        output.append(f"Style: {style}")
        output.append(f"{'='*80}\n")
        output.append(result)
        output.append(f"\n{'='*80}")
        output.append("NOTE: This is a draft. Review and edit as needed.")
        output.append(f"{'='*80}")

        print(f"[OK] {section_type} section drafted")
        return "\n".join(output)

    except Exception as e:
        return f"[ERROR] Error writing section: {str(e)}"


async def main():
    
    web_researcher: Agent = Agent(
        name="Web Research Specialist",
        instructions="""Search and download research papers. 
Default: Download 10 papers (5 Semantic Scholar + 5 Google Scholar), rank, keep top 5, delete rest.
If user says "N articles": Download N×2, rank, keep N, delete rest.

Steps: 1) Search both sources 2) Download PDFs 3) Read with read_pdf 4) Rank by relevance 5) Delete extras with delete_pdf
Always report kept files with locations: C:\\Users\\DELL\\Downloads\\filename.pdf""",
        handoff_description="Downloads, ranks, and manages research papers",
        model=llm_model,
        tools=[semantic_scholar_search, google_scholar_search, read_pdf, download_pdf, batch_download_pdfs, delete_pdf]
    )

    reader: Agent = Agent(
        name="Document Reader",
        instructions="""Read documents and extract text with PAGE NUMBERS for citations.
Tools: read_pdf (PDFs/scanned), read_word (Word), read_pptx (PowerPoint), read_image (images), extract_text_from_audio (audio).
Always include page/slide numbers. Return complete text.""",
        handoff_description="Reads documents with page numbers",
        model=llm_model,
        tools=[read_pdf, read_word, read_pptx, read_folder, list_files_in_folder, extract_text_from_audio, read_image]
    )

    download_deleater: Agent = Agent(
        name="File Manager",
        instructions="""Download PDFs from URLs or delete files. Use confirm=True when deleting. Report file location and size.""",
        handoff_description="Downloads/deletes files",
        model=llm_model,
        tools=[download_pdf, delete_pdf]
    )

    output_form: Agent = Agent(
        name="Output Generator",
        instructions="""Create Word/PDF documents or voice output. Use create_word_file, create_pdf, or voice_output tools. Save to Downloads folder.""",
        handoff_description="Creates documents and voice output",
        model=llm_model,
        tools=[create_word_file, create_pdf, voice_output]
    )


    head_agent: Agent = Agent(
        name="Research Assistant",
        instructions="""You are a research assistant. ONLY answer from documents/papers - NEVER from your own knowledge.

WORKFLOW:
1. Question without docs → Use web_research_agent to download papers, then read and cite
2. Question with docs → Use file_reader_agent to read, extract evidence with page numbers
3. Download request → Use web_research_agent, it downloads, ranks, keeps best papers
4. Create output → Use output_generator_agent for Word/PDF/voice

MUST DO: Always cite sources [filename.pdf, Page X]. Never guess or assume.""",
        model=llm_model,
        tools=[
            web_researcher.as_tool(
                tool_name="web_research_agent",
                tool_description="Downloads research papers from web, ranks by relevance, keeps best ones"
            ),
            reader.as_tool(
                tool_name="file_reader_agent", 
                tool_description="Reads PDF/Word/audio/images with page numbers"
            ),
            download_deleater.as_tool(
                tool_name="file_manager_agent",
                tool_description="Downloads specific URLs or deletes files"
            ),
            output_form.as_tool(
                tool_name="output_generator_agent",
                tool_description="Creates Word/PDF documents or voice output"
            )
        ]
    )
    
    # s_agent: Agent = Agent(
    #     name="SUMMARIZER",
    #     instructions="write very short summary of important things in text straight forward",
    #     model=llmm_model,
    # )
    
    with trace("agent"):
        convo: list[TResponseInputItem] = []
        while True:
            user_input = input("user: ")
            convo.append({"content": user_input, "role": "user"})
            result = ""

            maine = Runner.run_streamed(head_agent, input=convo)
            async for event in maine.stream_events():
                # We'll ignore the raw responses event deltas
                if event.type == "raw_response_event":
                    continue
                # When the agent updates, print that
                elif event.type == "agent_updated_stream_event":
                    print(f"Agent updated: {event.new_agent.name}")
                    continue
                # When items are generated, print them
                elif event.type == "run_item_stream_event":
                    if event.item.type == "tool_call_item":
                        print("-- Tool was called")
                    elif event.item.type == "tool_call_output_item":
                        print(f"-- Tool output: {event.item.output}")
                    elif event.item.type == "message_output_item":
                        result = ItemHelpers.text_message_output(event.item)
                        print(f"-- Message output:\n {result}")
                    else:
                        pass  # Ignore other event types

            print("=== Run complete ===")
            if result:
                # result = "summarize this: " + result
                # summery_agent = await Runner.run(s_agent, input=result)
                # summary = f"Agent Response: {summery_agent.final_output}"
                convo.append({"content": result, "role": "assistant"})
                # print(convo)

 
if __name__ == "__main__":
    asyncio.run(main())