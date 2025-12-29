"""
File Tools Module

Tools for creating files and managing the file system.
"""

import os
from pathlib import Path
from typing import Optional


def create_word_file(content: str, file_name: str = "output.docx") -> str:
    """Create a Word document from text content."""
    try:
        from docx import Document
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Add title
        title = doc.add_heading('Generated Document', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add content paragraphs
        for paragraph in content.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())

        # Save document
        downloads_folder = Path(__file__).parent.parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        file_path = downloads_folder / file_name
        doc.save(file_path)

        file_size = os.path.getsize(file_path) / (1024 * 1024)
        return f"[OK] Created Word document: {file_name} ({file_size:.2f} MB)\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: /api/files/download/{file_name}"

    except ImportError:
        return "[ERROR] python-docx not installed. Install with: pip install python-docx"
    except Exception as e:
        return f"[ERROR] Failed to create Word document: {str(e)}"


def create_pdf(content: str, file_name: str = "output.pdf") -> str:
    """Create a PDF from text content."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch

        downloads_folder = Path(__file__).parent.parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        file_path = downloads_folder / file_name

        doc = SimpleDocTemplate(str(file_path), pagesize=letter)
        styles = getSampleStyleSheet()

        # Custom style for body text
        body_style = ParagraphStyle(
            'BodyText',
            parent=styles['Normal'],
            fontSize=11,
            leading=14,
            spaceAfter=12
        )

        # Build document
        story = []
        story.append(Paragraph("Generated Document", styles['Title']))
        story.append(Spacer(1, 0.2 * inch))

        for paragraph in content.split('\n\n'):
            if paragraph.strip():
                # Escape HTML characters
                escaped = paragraph.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(escaped, body_style))
                story.append(Spacer(1, 0.1 * inch))

        doc.build(story)

        file_size = os.path.getsize(file_path) / (1024 * 1024)
        return f"[OK] Created PDF: {file_name} ({file_size:.2f} MB)\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: /api/files/download/{file_name}"

    except ImportError:
        return "[ERROR] reportlab not installed. Install with: pip install reportlab"
    except Exception as e:
        return f"[ERROR] Failed to create PDF: {str(e)}"


def create_pptx(content: str, file_name: str = "output.pptx", theme: str = "professional") -> str:
    """Create a PowerPoint presentation from text content."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        downloads_folder = Path(__file__).parent.parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        file_path = downloads_folder / file_name

        prs = Presentation()

        # Theme colors
        themes = {
            "professional": {
                "primary": RGBColor(0, 51, 102),  # Dark blue
                "secondary": RGBColor(51, 153, 255),  # Light blue
                "background": RGBColor(255, 255, 255)
            },
            "modern": {
                "primary": RGBColor(45, 45, 45),  # Dark gray
                "secondary": RGBColor(100, 100, 100),
                "background": RGBColor(30, 30, 30)
            },
            "elegant": {
                "primary": RGBColor(75, 0, 130),  # Purple
                "secondary": RGBColor(147, 112, 219),
                "background": RGBColor(255, 255, 255)
            },
            "nature": {
                "primary": RGBColor(34, 139, 34),  # Green
                "secondary": RGBColor(144, 238, 144),
                "background": RGBColor(255, 255, 255)
            },
            "warm": {
                "primary": RGBColor(204, 102, 0),  # Orange
                "secondary": RGBColor(255, 178, 102),
                "background": RGBColor(255, 255, 255)
            }
        }

        colors = themes.get(theme, themes["professional"])

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Research Presentation"
        subtitle.text = "Generated by Research Agent"

        # Content slides
        slide_layout = prs.slide_layouts[1]  # Title and Content

        # Split content into sections
        sections = content.split('\n\n')
        current_slide = None
        bullet_points = []

        for section in sections:
            if section.strip():
                # Create new slide every 5 bullet points or when section is long
                if len(bullet_points) >= 5 or len(section) > 500:
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = "Content"

                    tf = slide.placeholders[1].text_frame
                    tf.text = bullet_points[0]

                    for point in bullet_points[1:]:
                        p = tf.add_paragraph()
                        p.text = point

                    bullet_points = []

                bullet_points.append(section.strip()[:200])

        # Add remaining bullet points
        if bullet_points:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Content"

            tf = slide.placeholders[1].text_frame
            if bullet_points:
                tf.text = bullet_points[0]
                for point in bullet_points[1:]:
                    p = tf.add_paragraph()
                    p.text = point

        prs.save(file_path)

        file_size = os.path.getsize(file_path) / (1024 * 1024)
        slide_count = len(prs.slides)
        return f"[OK] Created PowerPoint: {file_name} ({file_size:.2f} MB, {slide_count} slides)\n[FILE]: {file_name}\n[DOWNLOAD_LINK]: /api/files/download/{file_name}"

    except ImportError:
        return "[ERROR] python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"[ERROR] Failed to create PowerPoint: {str(e)}"


def voice_output(text: str, filename: str = "output.wav") -> str:
    """Convert text to speech and save as audio file."""
    try:
        import edge_tts

        downloads_folder = Path(__file__).parent.parent / "downloads"
        downloads_folder.mkdir(exist_ok=True)
        file_path = downloads_folder / filename

        # Use edge-tts for free text-to-speech
        import asyncio

        async def generate():
            communicate = edge_tts.Communicate(text, "en-US-GuyNeural")
            await communicate.save(str(file_path))

        asyncio.run(generate())

        file_size = os.path.getsize(file_path) / (1024 * 1024)
        return f"[OK] Created audio: {filename} ({file_size:.2f} MB)\n[FILE]: {filename}\n[DOWNLOAD_LINK]: /api/files/download/{filename}"

    except ImportError:
        return "[ERROR] edge-tts not installed. Install with: pip install edge-tts"
    except Exception as e:
        return f"[ERROR] Failed to create audio: {str(e)}"


def read_folder(folder_path: str) -> str:
    """Read contents of a folder recursively."""
    try:
        from pathlib import Path

        path = Path(folder_path)
        if not path.exists():
            return f"[ERROR] Folder not found: {folder_path}"

        if not path.is_dir():
            return f"[ERROR] Not a folder: {folder_path}"

        results = []
        results.append(f"Folder: {folder_path}")
        results.append("=" * 80)

        for item in sorted(path.rglob("*")):
            if item.is_file():
                size = item.stat().st_size
                size_str = f"{size / 1024:.1f} KB" if size < 1024 * 1024 else f"{size / (1024 * 1024):.1f} MB"
                results.append(f"[FILE] {item.relative_to(path)} ({size_str})")
            elif item.is_dir():
                results.append(f"[DIR] {item.name}/")

        return "\n".join(results)

    except Exception as e:
        return f"[ERROR] Failed to read folder: {str(e)}"


def list_files_in_folder(folder_path: str, file_type: Optional[str] = None) -> str:
    """List files in a folder, optionally filtered by type."""
    try:
        from pathlib import Path

        path = Path(folder_path)
        if not path.exists():
            return f"[ERROR] Folder not found: {folder_path}"

        files = []
        for item in path.iterdir():
            if item.is_file():
                if file_type is None or item.suffix.lower() == file_type.lower():
                    files.append(item.name)

        if file_type:
            return f"Files in {folder_path} ({file_type}):\n" + "\n".join(sorted(files))
        return f"Files in {folder_path}:\n" + "\n".join(sorted(files))

    except Exception as e:
        return f"[ERROR] Failed to list files: {str(e)}"


def delete_file(filename: str) -> str:
    """Delete a file from downloads folder."""
    try:
        downloads_folder = Path(__file__).parent.parent / "downloads"
        file_path = downloads_folder / filename
        if file_path.exists():
            os.remove(file_path)
            return f"[OK] Deleted: {filename}"
        return f"[ERROR] File not found: {filename}"
    except Exception as e:
        return f"[ERROR] {str(e)}"
